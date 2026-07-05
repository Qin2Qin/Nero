#!/usr/bin/env python3
"""Build the Nero 6-slide pitch deck as a .pptx (imports cleanly into Google Slides)."""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SS = "/private/tmp/claude-501/-Users-user-Documents-Xero-Hackathon/b2310cfc-4915-404b-8c4e-16725db7157c/scratchpad"
OUT = "/Users/user/Documents/Xero Hackathon/docs/nero-pitch.pptx"
os.makedirs(os.path.dirname(OUT), exist_ok=True)

# ---- palette (matches the product's dark theme) ----
BG      = RGBColor(0x0A, 0x0E, 0x1A)
CARD    = RGBColor(0x12, 0x18, 0x28)
WHITE   = RGBColor(0xF4, 0xF6, 0xFB)
MUTED   = RGBColor(0x9C, 0xA8, 0xBD)
PURPLE  = RGBColor(0x8B, 0x5C, 0xF6)
TEAL    = RGBColor(0x2D, 0xD4, 0xBF)
RED     = RGBColor(0xF8, 0x71, 0x71)

# ---- crop tall screenshots + compress (keeps the .pptx light for upload) ----
def prep(src, dst, keep_frac=1.0, max_w=1200, q=88):
    im = Image.open(src).convert("RGB")
    w, h = im.size
    if keep_frac < 1.0:
        im = im.crop((0, 0, w, int(h * keep_frac))); w, h = im.size
    if w > max_w:
        im = im.resize((max_w, int(h * max_w / w)), Image.LANCZOS)
    im.save(dst, "JPEG", quality=q)
    return dst

dash_c  = prep(f"{SS}/deck-01-dashboard.png", f"{SS}/crop-dash.jpg", 0.56)
act_c   = prep(f"{SS}/deck-02-actions.png",  f"{SS}/crop-act.jpg",  0.66)
payers  = prep(f"{SS}/deck-03-payers.png",   f"{SS}/crop-payers.jpg", 1.0, max_w=1280, q=88)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = BG; r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s

def box(s, l, t, w, h, color):
    b = s.shapes.add_shape(5, Inches(l), Inches(t), Inches(w), Inches(h))
    b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()
    b.shadow.inherit = False
    return b

def text(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=1.0):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = sp
        if isinstance(para, dict):
            para = [para]
        for run in para:
            r = p.add_run(); r.text = run["t"]
            r.font.size = Pt(run.get("s", 18)); r.font.bold = run.get("b", False)
            r.font.color.rgb = run.get("c", WHITE); r.font.name = "Arial"
            if run.get("space_before"): p.space_before = Pt(run["space_before"])
        if isinstance(para, list) and para and para[0].get("space_before"):
            p.space_before = Pt(para[0]["space_before"])
    return tb

def pic_fit(s, path, l, t, maxw, maxh, align_right=False):
    im = Image.open(path); iw, ih = im.size; ar = iw / ih
    w = maxw; h = w / ar
    if h > maxh:
        h = maxh; w = h * ar
    left = l + (maxw - w) if align_right else l + (maxw - w) / 2
    p = s.shapes.add_picture(path, Inches(left), Inches(t + (maxh - h) / 2), Inches(w), Inches(h))
    p.line.color.rgb = RGBColor(0x2A, 0x33, 0x48); p.line.width = Pt(1)
    return p

def kicker(s, txt, l=0.9, t=0.55):
    text(s, l, t, 8, 0.4, [[{"t": txt, "s": 13, "b": True, "c": PURPLE}]])

# ---------- Slide 1 — Title ----------
s = slide()
box(s, 0.9, 2.35, 0.16, 2.7, PURPLE)
text(s, 1.3, 2.15, 11, 1.4, [[{"t": "Nero", "s": 66, "b": True, "c": WHITE}]])
text(s, 1.32, 3.45, 11, 0.7, [[{"t": "The Cash Flow Accelerator", "s": 28, "b": True, "c": TEAL}]])
text(s, 1.32, 4.35, 10.4, 1.4, [[
    {"t": "Your ledger tells you what you're owed. ", "s": 20, "c": MUTED},
    {"t": "Nero tells you when the cash actually arrives — and acts to bring it forward.", "s": 20, "c": WHITE},
]], sp=1.15)
text(s, 1.32, 6.5, 11, 0.5, [[{"t": "Bounty 03  ·  Xero App & Agent Hackathon  ·  an AR-intelligence agent for Xero", "s": 13, "c": MUTED}]])

# ---------- Slide 2 — Problem ----------
s = slide()
kicker(s, "THE PROBLEM")
text(s, 0.9, 1.05, 11.5, 1.0, [[{"t": "Charts track money. They don't get you paid.", "s": 34, "b": True, "c": WHITE}]])
text(s, 0.9, 2.15, 11.5, 1.1, [[
    {"t": "Due dates are fiction. ", "s": 20, "b": True, "c": WHITE},
    {"t": "A small business sees what's ", "s": 20, "c": MUTED},
    {"t": "due", "s": 20, "b": True, "c": WHITE},
    {"t": " — not when cash ", "s": 20, "c": MUTED},
    {"t": "actually", "s": 20, "b": True, "c": WHITE},
    {"t": " lands. So they hit cash crunches they could have prevented.", "s": 20, "c": MUTED},
]], sp=1.15)
cards = [("Due next 30 days", "£18,000", MUTED, "what the ledger shows"),
         ("Actually arriving", "£9,000", PURPLE, "once you price in late payers"),
         ("The gap", "£9,000", RED, "invisible until it's a crisis")]
for i, (lab, val, col, sub) in enumerate(cards):
    x = 0.9 + i * 4.05
    box(s, x, 3.75, 3.75, 2.7, CARD)
    text(s, x + 0.35, 4.05, 3.1, 0.4, [[{"t": lab.upper(), "s": 12, "b": True, "c": MUTED}]])
    text(s, x + 0.35, 4.5, 3.1, 1.0, [[{"t": val, "s": 46, "b": True, "c": col}]])
    text(s, x + 0.35, 5.75, 3.1, 0.6, [[{"t": sub, "s": 13, "c": MUTED}]])

# ---------- Slide 3 — Product (dashboard) ----------
s = slide()
kicker(s, "THE PRODUCT  ·  TRUTH → FORECAST")
text(s, 0.9, 1.05, 6.1, 1.1, [[{"t": "See the real cash picture", "s": 32, "b": True, "c": WHITE}]])
for i, (h, b) in enumerate([
    ("Reads your Xero ledger", "Contacts, invoices and payment history — no manual entry."),
    ("Predicts real payment dates", "Each customer's own behaviour, not the due date."),
    ("Forecasts against a floor", "Flags the exact week cash dips below your minimum."),
]):
    y = 2.2 + i * 1.35
    box(s, 0.9, y, 0.1, 1.05, TEAL)
    text(s, 1.2, y, 5.6, 1.2, [
        [{"t": h, "s": 18, "b": True, "c": WHITE}],
        [{"t": b, "s": 14, "c": MUTED, "space_before": 3}],
    ], sp=1.05)
pic_fit(s, dash_c, 7.2, 1.0, 5.5, 6.1, align_right=True)

# ---------- Slide 4 — Agent (actions) ----------
s = slide()
kicker(s, "THE AGENT  ·  ACTION, WITH YOUR APPROVAL")
text(s, 0.9, 1.05, 6.1, 1.1, [[{"t": "An agent that acts —\nyou stay in control", "s": 30, "b": True, "c": WHITE}]], sp=1.0)
text(s, 0.9, 2.7, 5.9, 2.0, [[
    {"t": "Nero proposes the specific moves — ", "s": 16, "c": MUTED},
    {"t": "chase this invoice, ask a deposit there, tighten terms", "s": 16, "b": True, "c": WHITE},
    {"t": " — each tagged with the £ it brings forward. ", "s": 16, "c": MUTED},
    {"t": "You approve; nothing sends automatically.", "s": 16, "b": True, "c": TEAL},
]], sp=1.15)
box(s, 0.9, 5.2, 5.9, 1.4, CARD)
text(s, 1.2, 5.45, 5.4, 1.0, [
    [{"t": "£13,000", "s": 30, "b": True, "c": TEAL}, {"t": "  ready to bring forward", "s": 16, "c": WHITE}],
    [{"t": "~12 days sooner, across 7 approval-gated actions", "s": 13, "c": MUTED, "space_before": 4}],
])
pic_fit(s, act_c, 7.1, 0.95, 5.6, 6.2, align_right=True)

# ---------- Slide 5 — Payer intelligence (payers) ----------
s = slide()
kicker(s, "PAYER INTELLIGENCE")
text(s, 0.9, 1.05, 11.5, 1.0, [[{"t": "Know exactly who's slowing your cash", "s": 30, "b": True, "c": WHITE}]])
text(s, 0.9, 2.0, 11.5, 0.6, [[
    {"t": "Every customer graded on their real payment behaviour, ranked by cash at risk — ", "s": 15, "c": MUTED},
    {"t": "who to chase first.", "s": 15, "b": True, "c": WHITE},
]])
pic_fit(s, payers, 0.9, 2.65, 11.5, 4.5)

# ---------- Slide 6 — Why Nero wins ----------
s = slide()
kicker(s, "WHY NERO WINS")
text(s, 0.9, 1.05, 11.5, 0.9, [[{"t": "Accountable by design, built the Xero way", "s": 30, "b": True, "c": WHITE}]])
cols = [
    ("Built the Xero way", TEAL, [
        "Live OAuth to a Xero organisation",
        "Reads contacts, invoices, payments +",
        "   online invoice links (Accounting API)",
        "Approved reminders write a history",
        "   note back to the invoice in Xero",
    ]),
    ("Accountable by design", PURPLE, [
        "Every action is approval-gated",
        "Reasoning shown, grounded in real stats",
        "Nothing sends automatically",
        "\"Cash Accelerated £\" = measurable outcome",
    ]),
    ("Validated demand", WHITE, [
        "Payment-performance report: 351 forum votes",
        "Deposits on invoices: 507 votes",
        "…both unbuilt by Xero today",
        "Chases + forecasts in one loop (nobody else)",
    ]),
]
for i, (title, col, items) in enumerate(cols):
    x = 0.9 + i * 4.05
    box(s, x, 2.2, 3.8, 3.5, CARD)
    text(s, x + 0.3, 2.45, 3.2, 0.5, [[{"t": title, "s": 16, "b": True, "c": col}]])
    text(s, x + 0.3, 3.15, 3.25, 2.4,
         [[{"t": "•  " + it, "s": 12.5, "c": MUTED}] for it in items], sp=1.25)
text(s, 0.9, 6.05, 11.5, 0.9, [[
    {"t": "What's next:  ", "s": 14, "b": True, "c": WHITE},
    {"t": "outcome learning  ·  SMS / multi-channel reminders  ·  multi-client advisor view", "s": 14, "c": MUTED},
]])
box(s, 0.9, 6.85, 11.53, 0.06, PURPLE)

prs.save(OUT)
print("SAVED", OUT, os.path.getsize(OUT), "bytes,", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
